"""Build Promptly pitch deck as a .pptx file."""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color tokens ──────────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x14, 0x14, 0x14)
INK2 = RGBColor(0x1E, 0x1E, 0x22)
MUTED = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE1)
VIOLET = RGBColor(0x7C, 0x5C, 0xFF)
PINK = RGBColor(0xFF, 0x7A, 0xC6)
GREEN = RGBColor(0x5C, 0xFF, 0xB1)
AMBER = RGBColor(0xFF, 0xB8, 0x5C)
FG = RGBColor(0xED, 0xED, 0xED)
FG2 = RGBColor(0xB5, 0xB5, 0xBA)
FG3 = RGBColor(0x8A, 0x8A, 0x90)

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # fully blank layout


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=None, border_color=None, border_pt=0.75):
    """Add a rectangle shape."""
    from pptx.util import Pt as Pt_

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.line.width = Pt_(border_pt) if border_color else 0
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def txt(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    bold=False,
    italic=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    font_name="Calibri",
    wrap=True,
):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def para(
    tf,
    text,
    size=14,
    bold=False,
    italic=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    font_name="Calibri",
    space_before=0,
):
    """Add a paragraph to an existing text frame."""
    from pptx.util import Pt as Pt_

    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt_(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt_(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def eyebrow(slide, text, x, y):
    """Small violet uppercase mono label."""
    txt(
        slide,
        text,
        x,
        y,
        10,
        0.3,
        size=9,
        bold=True,
        color=VIOLET,
        font_name="Courier New",
    )


def divider(slide, y, x=0.5, w=12.33):
    """Thin horizontal line."""
    ln = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.01))
    ln.fill.solid()
    ln.fill.fore_color.rgb = LIGHT_GRAY
    ln.line.fill.background()


def stat_block(slide, number, label, x, y, dark=False):
    """A large number + small label block."""
    nc = WHITE if dark else VIOLET
    lc = FG2 if dark else MUTED
    txt(
        slide,
        number,
        x,
        y,
        2.8,
        0.85,
        size=44,
        bold=True,
        color=nc,
        font_name="Georgia",
        align=PP_ALIGN.CENTER,
    )
    txt(slide, label, x, y + 0.85, 2.8, 0.5, size=10, color=lc, align=PP_ALIGN.CENTER)


def agent_chip(slide, letter, name, color, x, y):
    """Colored chip with letter + name."""
    box(slide, x, y, 0.35, 0.35, fill_color=color)
    txt(
        slide,
        letter,
        x + 0.02,
        y + 0.03,
        0.31,
        0.3,
        size=13,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    txt(slide, name, x + 0.42, y + 0.06, 1.4, 0.25, size=10, bold=False, color=INK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, INK2)

# Logo area
box(s, 0.5, 0.4, 0.3, 0.3, fill_color=VIOLET)
txt(s, "promptly", 0.88, 0.42, 2, 0.28, size=14, bold=True, color=WHITE)

# Headline
txt(
    s,
    "Your prompt,",
    0.5,
    1.4,
    12,
    1.0,
    size=52,
    bold=False,
    italic=False,
    color=WHITE,
    font_name="Georgia",
)
txt(
    s,
    "sharpened",
    0.5,
    2.3,
    12,
    1.0,
    size=52,
    bold=False,
    italic=True,
    color=VIOLET,
    font_name="Georgia",
)
txt(
    s,
    "by a society of minds.",
    0.5,
    3.2,
    12,
    1.0,
    size=52,
    bold=False,
    color=WHITE,
    font_name="Georgia",
)

# Subtext
txt(
    s,
    "Promptly — AI prompt optimization platform",
    0.5,
    4.5,
    10,
    0.4,
    size=14,
    color=FG3,
    font_name="Courier New",
)

# Council dots strip
for i, (letter, color) in enumerate(
    [("A", VIOLET), ("C", PINK), ("O", GREEN), ("S", AMBER)]
):
    cx = 0.5 + i * 0.55
    box(s, cx, 6.5, 0.32, 0.32, fill_color=color)
    txt(
        s,
        letter,
        cx + 0.02,
        6.52,
        0.28,
        0.28,
        size=12,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

txt(
    s,
    "Analytical  ·  Creative  ·  Concise  ·  Structured",
    2.8,
    6.52,
    8,
    0.3,
    size=10,
    color=FG3,
    font_name="Courier New",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, WHITE)

eyebrow(s, "THE PROBLEM", 0.5, 0.4)
txt(
    s,
    "Every team shipping LLM features\nloses weeks to prompt tuning.",
    0.5,
    0.85,
    12,
    1.5,
    size=36,
    bold=False,
    color=INK,
    font_name="Georgia",
)

divider(s, 2.5)

# Three pain-point cards
for i, (title, body) in enumerate(
    [
        (
            "Manual iteration",
            "Hours → days → weeks of back-and-forth. Every model update restarts the clock.",
        ),
        (
            "Silent breakage",
            "Prompts break across model versions with no warning. You find out in production.",
        ),
        (
            "No history",
            "No versioning, no rollback, no audit trail. Every prompt is a one-off.",
        ),
    ]
):
    cx = 0.5 + i * 4.1
    box(s, cx, 2.7, 3.8, 2.5, fill_color=None, border_color=LIGHT_GRAY)
    txt(s, title, cx + 0.2, 2.85, 3.4, 0.4, size=13, bold=True, color=INK)
    txt(s, body, cx + 0.2, 3.3, 3.4, 1.7, size=11, color=MUTED)

# Closing line
divider(s, 5.4)
txt(
    s,
    '"You blame the model. It\'s never the model."',
    0.5,
    5.55,
    12,
    0.5,
    size=16,
    italic=True,
    color=VIOLET,
    font_name="Georgia",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — COST OF A BAD PROMPT
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, WHITE)

eyebrow(s, "REAL COST", 0.5, 0.4)
txt(
    s,
    "A bad prompt isn't a bug.\nIt's a revenue leak.",
    0.5,
    0.85,
    10,
    1.2,
    size=38,
    bold=False,
    color=INK,
    font_name="Georgia",
)

divider(s, 2.2)

# Three stat cards
for i, (num, label, source) in enumerate(
    [
        (
            "60%",
            "of AI feature failures trace back\nto prompt quality",
            "Statista, 2025",
        ),
        (
            "340%",
            "ROI from optimized prompts\nvs. unoptimized baselines",
            "Deloitte, 2025",
        ),
        (
            "Weeks",
            "of engineer time lost per model\nmigration — repeated every release",
            "Industry avg.",
        ),
    ]
):
    cx = 0.5 + i * 4.1
    box(s, cx, 2.4, 3.8, 3.2, fill_color=None, border_color=LIGHT_GRAY)
    txt(
        s,
        num,
        cx + 0.2,
        2.6,
        3.4,
        0.9,
        size=44,
        bold=True,
        color=VIOLET,
        font_name="Georgia",
    )
    txt(s, label, cx + 0.2, 3.5, 3.4, 1.0, size=11, color=INK)
    txt(
        s, source, cx + 0.2, 5.2, 3.4, 0.3, size=9, color=MUTED, font_name="Courier New"
    )


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, WHITE)

eyebrow(s, "MARKET OPPORTUNITY", 0.5, 0.4)
txt(
    s,
    "Every AI product needs this.\nThe market knows it.",
    0.5,
    0.85,
    10,
    1.2,
    size=38,
    bold=False,
    color=INK,
    font_name="Georgia",
)

divider(s, 2.2)

for i, (num, label) in enumerate(
    [
        ("$1.7B", "Prompt optimization market\nsize in 2024"),
        ("$15.2B", "Projected market size\nby 2033"),
        ("26%", "CAGR — one of the fastest\ngrowing AI sub-markets"),
    ]
):
    cx = 0.5 + i * 4.1
    box(s, cx, 2.4, 3.8, 2.8, fill_color=None, border_color=LIGHT_GRAY)
    txt(
        s,
        num,
        cx + 0.2,
        2.6,
        3.4,
        0.9,
        size=44,
        bold=True,
        color=VIOLET,
        font_name="Georgia",
    )
    txt(s, label, cx + 0.2, 3.55, 3.4, 0.9, size=11, color=INK)

txt(
    s,
    "Source: Growth Market Reports, 2025",
    0.5,
    5.45,
    8,
    0.3,
    size=9,
    color=MUTED,
    font_name="Courier New",
)
divider(s, 5.85)
txt(
    s,
    '"This is infrastructure, not tooling."',
    0.5,
    6.0,
    12,
    0.4,
    size=15,
    italic=True,
    color=VIOLET,
    font_name="Georgia",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE INSIGHT: SOCIETY OF MIND
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, INK2)

eyebrow(s, "THE INSIGHT", 0.5, 0.4)

# Large quote mark
txt(s, "“", 0.5, 0.8, 1.5, 1.5, size=96, color=VIOLET, font_name="Georgia")

txt(
    s,
    "The mind is a society of tiny agents,\neach mindless by itself.\nIntelligence is what emerges\nwhen they interact.",
    1.2,
    1.1,
    10.5,
    2.8,
    size=30,
    italic=True,
    color=FG,
    font_name="Georgia",
)

txt(
    s,
    "— Marvin Minsky  ·  Society of Mind, 1986",
    1.2,
    3.95,
    10,
    0.4,
    size=11,
    color=FG3,
    font_name="Courier New",
)

divider(s, 4.55)

txt(
    s,
    "Minsky proved that intelligence is emergent, not monolithic. In 1986, he described\nthe architecture we built — now applied to prompt engineering.",
    0.5,
    4.75,
    12,
    0.9,
    size=13,
    color=FG2,
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — INTRODUCING PROMPTLY
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, WHITE)

eyebrow(s, "THE SOLUTION", 0.5, 0.4)
txt(
    s,
    "A society of four specialist agents —\nproposing, reviewing, synthesizing.",
    0.5,
    0.85,
    12,
    1.3,
    size=34,
    bold=False,
    color=INK,
    font_name="Georgia",
)

txt(
    s,
    "Paste any prompt. Get back the best version none of the agents could write alone.",
    0.5,
    2.3,
    11,
    0.5,
    size=14,
    color=MUTED,
)

divider(s, 3.0)

# Agent chips
for i, (letter, name, color) in enumerate(
    [
        ("A", "Analytical", VIOLET),
        ("C", "Creative", PINK),
        ("O", "Concise", GREEN),
        ("S", "Structured", AMBER),
    ]
):
    agent_chip(s, letter, name, color, 0.5 + i * 3.1, 3.2)

divider(s, 3.85)

# Pipeline
pipeline = ["INPUT", "PROPOSE", "CRITIQUE", "SYNTHESIZE", "OUTPUT"]
for i, step in enumerate(pipeline):
    cx = 0.5 + i * 2.5
    is_key = step in ("PROPOSE", "CRITIQUE", "SYNTHESIZE")
    fc = VIOLET if is_key else LIGHT_GRAY
    tc = WHITE if is_key else INK
    box(s, cx, 4.1, 2.1, 0.55, fill_color=fc, border_color=None)
    txt(
        s,
        step,
        cx,
        4.18,
        2.1,
        0.4,
        size=11,
        bold=True,
        color=tc,
        align=PP_ALIGN.CENTER,
        font_name="Courier New",
    )
    if i < 4:
        txt(
            s,
            "→",
            cx + 2.1,
            4.22,
            0.35,
            0.35,
            size=14,
            color=VIOLET,
            align=PP_ALIGN.CENTER,
        )


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — HOW IT WORKS: PROPOSE
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, WHITE)

eyebrow(s, "01  ·  PROPOSE", 0.5, 0.4)
txt(
    s,
    "Four minds. One prompt. Parallel.",
    0.5,
    0.85,
    12,
    0.8,
    size=38,
    bold=False,
    color=INK,
    font_name="Georgia",
)

divider(s, 1.8)

agents = [
    (
        "A",
        "Analytical",
        VIOLET,
        "Adds constraints, output schema, and precision guardrails.",
    ),
    ("C", "Creative", PINK, "Adds persona, worked examples, and voice directives."),
    ("O", "Concise", GREEN, "Strips every filler word. Maximises signal density."),
    ("S", "Structured", AMBER, "Adds logical decomposition and output schemas."),
]
for i, (letter, name, color, desc) in enumerate(agents):
    col = i % 2
    row = i // 2
    cx = 0.5 + col * 6.2
    cy = 2.0 + row * 2.3
    box(s, cx, cy, 5.9, 2.0, fill_color=None, border_color=LIGHT_GRAY)
    box(s, cx + 0.2, cy + 0.25, 0.4, 0.4, fill_color=color)
    txt(
        s,
        letter,
        cx + 0.22,
        cy + 0.27,
        0.36,
        0.36,
        size=14,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    txt(
        s,
        name.upper(),
        cx + 0.75,
        cy + 0.3,
        4.5,
        0.3,
        size=10,
        bold=True,
        color=INK,
        font_name="Courier New",
    )
    txt(s, desc, cx + 0.2, cy + 0.85, 5.4, 0.9, size=12, color=MUTED)

divider(s, 6.6)
txt(
    s,
    "No agent sees what the others write.",
    0.5,
    6.75,
    12,
    0.4,
    size=11,
    color=MUTED,
    font_name="Courier New",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — HOW IT WORKS: CRITIQUE
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, WHITE)

eyebrow(s, "02  ·  CRITIQUE", 0.5, 0.4)
txt(
    s,
    "Blind peer review. No favouritism.",
    0.5,
    0.85,
    12,
    0.8,
    size=38,
    bold=False,
    color=INK,
    font_name="Georgia",
)

divider(s, 1.8)

txt(
    s,
    "Each agent reviews the other three proposals — shown anonymously as A/B/C/D.",
    0.5,
    2.0,
    12,
    0.4,
    size=13,
    color=MUTED,
)

# Ranking table
box(s, 0.5, 2.6, 6.5, 2.8, fill_color=INK2)
rows = [
    ("critic A", "B > D > A > C"),
    ("critic B", "A > B > D > C"),
    ("critic C", "C > A > D > B"),
    ("critic D", "D > A > B > C"),
]
for i, (critic, ranking) in enumerate(rows):
    ry = 2.75 + i * 0.6
    txt(s, critic, 0.7, ry, 1.5, 0.5, size=11, color=FG3, font_name="Courier New")
    txt(
        s,
        ranking,
        2.2,
        ry,
        4.5,
        0.5,
        size=13,
        bold=True,
        color=FG,
        font_name="Courier New",
    )

# Right callout
txt(
    s,
    "Rankings aggregated\n→ consensus emerges",
    7.5,
    2.9,
    5.0,
    1.0,
    size=20,
    italic=True,
    color=VIOLET,
    font_name="Georgia",
)
txt(
    s,
    "No agent knows which proposal is theirs.\nNo model favours its own work.",
    7.5,
    4.1,
    5.0,
    0.8,
    size=12,
    color=MUTED,
)

divider(s, 5.7)
txt(
    s,
    "The result: a rankings consensus that emerges from four independent judgements.",
    0.5,
    5.85,
    12,
    0.4,
    size=12,
    italic=True,
    color=MUTED,
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — HOW IT WORKS: SYNTHESIZE
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, INK2)

eyebrow(s, "03  ·  SYNTHESIZE", 0.5, 0.4)
txt(
    s,
    "A chairman writes the final.",
    0.5,
    0.85,
    12,
    0.8,
    size=40,
    bold=False,
    color=WHITE,
    font_name="Georgia",
)

txt(
    s,
    "A separate synthesis model reads all four proposals and all four critique reports,\nthen extracts the strongest element from each.",
    0.5,
    1.8,
    12,
    0.7,
    size=13,
    color=FG2,
)

divider(s, 2.7)

extractions = [
    ("→  structure    from proposal D", AMBER),
    ("→  persona      from proposal B", PINK),
    ("→  constraints  from proposal A", VIOLET),
    ("→  brevity      from proposal C", GREEN),
]
for i, (line, color) in enumerate(extractions):
    txt(
        s,
        line,
        0.7,
        3.0 + i * 0.72,
        10,
        0.6,
        size=18,
        bold=False,
        color=color,
        font_name="Courier New",
    )

divider(s, 5.95)
txt(
    s,
    "One result. Genuinely emergent. None of them could write it alone.",
    0.5,
    6.1,
    12,
    0.5,
    size=16,
    italic=True,
    color=WHITE,
    font_name="Georgia",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, WHITE)

eyebrow(s, "RESULTS", 0.5, 0.4)
txt(
    s,
    "The numbers speak.",
    0.5,
    0.85,
    12,
    0.7,
    size=42,
    bold=False,
    color=INK,
    font_name="Georgia",
)

divider(s, 1.75)

stats = [
    ("3.8×", "average quality uplift\nvs. original prompt"),
    ("42s", "median end-to-end\noptimize time"),
    ("12,400", "prompts optimized\nthis month"),
    ("94%", "of users ship the synthesized\nresult unchanged"),
]
for i, (num, label) in enumerate(stats):
    cx = 0.5 + i * 3.1
    if i < 3:
        box(s, cx + 3.05, 1.95, 0.05, 4.5, fill_color=LIGHT_GRAY)
    txt(
        s,
        num,
        cx,
        2.1,
        3.0,
        1.1,
        size=52,
        bold=True,
        color=VIOLET,
        font_name="Georgia",
        align=PP_ALIGN.CENTER,
    )
    txt(s, label, cx, 3.25, 3.0, 0.8, size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — WHO IT'S FOR
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, WHITE)

eyebrow(s, "WHO IT'S FOR", 0.5, 0.4)
txt(
    s,
    "Built for people who ship with LLMs.",
    0.5,
    0.85,
    12,
    0.8,
    size=38,
    bold=False,
    color=INK,
    font_name="Georgia",
)

divider(s, 1.8)

audiences = [
    (
        "for product teams",
        "Stop losing a week to prompt tuning.",
        "Paste the prompt from your PRD. Ship the optimized one to staging by lunch. The council handles role, format, constraints, and edge cases automatically.",
    ),
    (
        "for engineers",
        "Treat prompts like code.",
        "Stable IDs, versioned history, diff view, rollback. Wire the API into your CI so every merge runs the health score against all ten quality dimensions.",
    ),
    (
        "for writers & ops",
        "Write in plain English.",
        "No prompt-engineering jargon. Say what you want; the society handles the rest. You'll never manually add a role, output schema, or guardrail again.",
    ),
]
for i, (eyebrow_t, headline, body) in enumerate(audiences):
    cx = 0.5 + i * 4.1
    box(s, cx, 2.05, 3.8, 4.5, fill_color=None, border_color=LIGHT_GRAY)
    txt(
        s,
        eyebrow_t.upper(),
        cx + 0.2,
        2.25,
        3.4,
        0.3,
        size=9,
        bold=True,
        color=VIOLET,
        font_name="Courier New",
    )
    txt(s, headline, cx + 0.2, 2.65, 3.4, 0.7, size=17, color=INK, font_name="Georgia")
    txt(s, body, cx + 0.2, 3.45, 3.4, 2.8, size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, WHITE)

eyebrow(s, "BUSINESS MODEL", 0.5, 0.4)
txt(
    s,
    "Pay for what you actually run.",
    0.5,
    0.85,
    12,
    0.7,
    size=40,
    bold=False,
    color=INK,
    font_name="Georgia",
)

divider(s, 1.75)

plans = [
    (
        "FREE",
        "$0",
        "/ forever",
        False,
        [
            "100 credits on signup",
            "10 credits / month refill",
            "3 saved prompt families",
            "Community support",
        ],
    ),
    (
        "PRO",
        "$29",
        "/ month",
        True,
        [
            "1,000 credits / month",
            "Unlimited prompt families",
            "API access + CLI",
            "Priority email support",
        ],
    ),
    (
        "TEAM",
        "$99",
        "/ month",
        False,
        [
            "5,000 credits pooled",
            "Up to 10 seats",
            "SSO + audit log",
            "Dedicated Slack channel",
        ],
    ),
]
for i, (plan, price, per, featured, features) in enumerate(plans):
    cx = 0.5 + i * 4.1
    card_bg = INK2 if featured else WHITE
    card_border = VIOLET if featured else LIGHT_GRAY
    box(
        s,
        cx,
        1.95,
        3.8,
        5.0,
        fill_color=card_bg,
        border_color=card_border,
        border_pt=1.5 if featured else 0.75,
    )

    txt(
        s,
        plan,
        cx + 0.2,
        2.15,
        3.4,
        0.35,
        size=10,
        bold=True,
        color=(FG3 if featured else MUTED),
        font_name="Courier New",
    )
    txt(
        s,
        price,
        cx + 0.2,
        2.55,
        2.2,
        0.8,
        size=40,
        bold=True,
        color=(WHITE if featured else INK),
        font_name="Georgia",
    )
    txt(s, per, cx + 2.3, 2.85, 1.4, 0.35, size=11, color=(FG3 if featured else MUTED))

    for j, feat in enumerate(features):
        fc = FG2 if featured else MUTED
        txt(s, f"✓  {feat}", cx + 0.2, 3.55 + j * 0.48, 3.4, 0.4, size=11, color=fc)

    if featured:
        txt(
            s,
            "MOST POPULAR",
            cx + 1.8,
            2.15,
            1.8,
            0.3,
            size=8,
            bold=True,
            color=VIOLET,
            font_name="Courier New",
        )

divider(s, 7.1)
txt(
    s,
    "1 optimize run = 10 credits  ·  1 health score = 5 credits  ·  unused credits roll 90 days",
    0.5,
    7.2,
    12,
    0.25,
    size=9,
    color=MUTED,
    font_name="Courier New",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — COMPETITIVE MOAT
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, WHITE)

eyebrow(s, "WHY WE WIN", 0.5, 0.4)
txt(
    s,
    "No one else runs four agents,\nblind-reviewing each other.",
    0.5,
    0.85,
    12,
    1.2,
    size=36,
    bold=False,
    color=INK,
    font_name="Georgia",
)

divider(s, 2.2)

# Table header
box(s, 0.5, 2.35, 8.5, 0.5, fill_color=INK2)
txt(
    s,
    "FEATURE",
    0.7,
    2.45,
    4.5,
    0.35,
    size=10,
    bold=True,
    color=FG3,
    font_name="Courier New",
)
txt(
    s,
    "PROMPTLY",
    5.3,
    2.45,
    1.8,
    0.35,
    size=10,
    bold=True,
    color=VIOLET,
    font_name="Courier New",
    align=PP_ALIGN.CENTER,
)
txt(
    s,
    "OTHERS",
    7.1,
    2.45,
    1.7,
    0.35,
    size=10,
    bold=True,
    color=FG3,
    font_name="Courier New",
    align=PP_ALIGN.CENTER,
)

rows: list[str] = [
    "Multi-agent peer review",
    "Emergent synthesis (not best-of-N)",
    "Stable prompt IDs + versioning",
    "Model-agnostic output",
    "Credit refund if no quality uplift",
]
for i, row in enumerate(rows):
    ry = 2.95 + i * 0.58
    rc = LIGHT_GRAY if i % 2 == 0 else WHITE
    box(s, 0.5, ry, 8.5, 0.55, fill_color=rc, border_color=None)
    box(s, 0.5, ry, 0.06, 0.55, fill_color=VIOLET)
    txt(s, row, 0.75, ry + 0.1, 4.4, 0.38, size=12, color=INK)
    txt(
        s,
        "✓",
        5.3,
        ry + 0.08,
        1.8,
        0.38,
        size=16,
        bold=True,
        color=VIOLET,
        align=PP_ALIGN.CENTER,
    )
    txt(
        s,
        "✗",
        7.1,
        ry + 0.08,
        1.7,
        0.38,
        size=16,
        bold=True,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

divider(s, 5.95)
txt(
    s,
    "Workflow lock-in via versioned prompt families. Data flywheel improves synthesis with every run.",
    0.5,
    6.1,
    12,
    0.4,
    size=11,
    italic=True,
    color=MUTED,
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CTA / CLOSE
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, INK2)

# Glowing dot + eyebrow
box(s, 0.5, 0.55, 0.12, 0.12, fill_color=VIOLET)
txt(
    s,
    "READY WHEN YOU ARE",
    0.72,
    0.52,
    8,
    0.3,
    size=10,
    bold=True,
    color=VIOLET,
    font_name="Courier New",
)

txt(
    s,
    "Let the society",
    0.5,
    1.3,
    12,
    0.9,
    size=54,
    bold=False,
    color=WHITE,
    font_name="Georgia",
)
txt(
    s,
    "sharpen",
    0.5,
    2.15,
    12,
    0.9,
    size=54,
    bold=False,
    italic=True,
    color=VIOLET,
    font_name="Georgia",
)
txt(
    s,
    "your prompt.",
    0.5,
    3.0,
    12,
    0.9,
    size=54,
    bold=False,
    color=WHITE,
    font_name="Georgia",
)

txt(
    s,
    "100 credits free  ·  No card required  ·  First council run under a minute.",
    0.5,
    4.15,
    12,
    0.45,
    size=14,
    color=FG2,
)

# CTA button visual
box(s, 0.5, 4.85, 3.2, 0.65, fill_color=VIOLET)
txt(
    s,
    "⚡  Try Promptly free",
    0.5,
    4.93,
    3.2,
    0.5,
    size=14,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)

divider(s, 6.1)
txt(
    s,
    "hey@promptly.dev  ·  promptly.dev",
    0.5,
    6.25,
    12,
    0.4,
    size=11,
    color=FG3,
    font_name="Courier New",
)

# Logo bottom center
box(s, 6.2, 6.9, 0.28, 0.28, fill_color=VIOLET)
txt(s, "promptly", 6.56, 6.92, 2, 0.25, size=12, bold=True, color=FG2)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
os.makedirs("/Volumes/External/promptly/docs/pitch-deck", exist_ok=True)
out = "/Volumes/External/promptly/docs/pitch-deck/promptly-pitch-deck.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
